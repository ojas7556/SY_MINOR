# app.py
import os
import json
import re
from io import BytesIO
from typing import List, Dict
from dotenv import load_dotenv
import streamlit as st
st.set_page_config(page_title="Virtual AI Tutor — Notes & Quiz", layout="wide", page_icon="🎓")
from fpdf import FPDF
from openai import OpenAI
from unidecode import unidecode   # transliterate unicode -> ascii (avoids latin-1 issues)
import requests
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from unidecode import unidecode   # transliterate unicode -> ascii (avoids latin-1 issues)

# -----------------------
# Config
# -----------------------
# -------- robust startup & secrets handling (replace existing top-of-file logic) ----------
import os
import traceback

# Optional: show full tracebacks in the app only while debugging.
# Remove or set DEBUG=False in production.
DEBUG = True

try:
    # Try imports that may fail on Cloud
    from dotenv import load_dotenv
    from openai import OpenAI
    # If you rely on load_dotenv locally, call it
    if os.getenv("STREAMLIT_RUNTIME") is None:  # naive check; adjust if you prefer
        load_dotenv()
except Exception:
    if DEBUG:
        st.error("Startup import error — see details below.")
        st.text(traceback.format_exc())
    raise

# Load key: prefer Streamlit secrets (Cloud) but fall back to environment variables (local)
OPENAI_API_KEY = None
try:
    OPENAI_API_KEY = st.secrets.get("OPENAI_API_KEY")  # safe .get() avoids KeyError
except Exception:
    # if st.secrets isn't available for any reason, ignore and fall back
    OPENAI_API_KEY = None

if not OPENAI_API_KEY:
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

MODEL_NAME = os.getenv("MODEL_NAME", "gpt-4o")

if not OPENAI_API_KEY:
    st.error(
        "OpenAI API key not found.\n\n"
        "• For Streamlit Cloud: Manage app → Settings → Secrets → add:\n"
        "    OPENAI_API_KEY = \"sk-...\"\n"
        "• For local dev: create .streamlit/secrets.toml or set env var OPENAI_API_KEY\n\n"
        "After adding the secret, redeploy or restart the app."
    )
    st.stop()

# create client
client = OpenAI(api_key=OPENAI_API_KEY)
# ---------------------------------------------------------------------------------------

# -----------------------
# History Helpers
# -----------------------
HISTORY_FILE = "history.json"

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return []
    return []

def save_history(record):
    history = load_history()
    history.insert(0, record)
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, indent=2)

# -----------------------
# DALL-E Image Generation
# -----------------------
def generate_image_with_dalle(prompt: str, size: str = "1024x1024") -> str:
    """Generate an image using OpenAI Image API. Returns image URL/data-URI or error message."""
    try:
        response = client.images.generate(
            model="gpt-image-1",
            prompt=prompt,
            size=size,
            n=1,
        )
        img_data = response.data[0]
        # Newer models may return b64_json instead of url
        if hasattr(img_data, 'url') and img_data.url:
            return img_data.url
        elif hasattr(img_data, 'b64_json') and img_data.b64_json:
            return f"data:image/png;base64,{img_data.b64_json}"
        else:
            return "Error generating image: No image data in response"
    except Exception as e:
        return f"Error generating image: {str(e)}"

def generate_multiple_images_with_dalle(topic: str, num_images: int = 3, context: str = "") -> List[str]:
    """Generate multiple images for different aspects of a topic."""
    image_urls = []
    
    # Create different prompts for different aspects
    aspects = [
        f"Overview diagram for {topic}",
        f"Detailed process flow for {topic}",
        f"Key concepts illustration for {topic}",
        f"Examples and applications of {topic}",
        f"Advanced topics in {topic}"
    ]
    
    for i in range(min(num_images, len(aspects))):
        prompt = create_image_prompt_from_topic(aspects[i], context)
        image_url = generate_image_with_dalle(prompt)
        if not image_url.startswith("Error"):
            image_urls.append(image_url)
        else:
            st.warning(f"Failed to generate image {i+1}: {image_url}")
    
    return image_urls


def extract_table_of_contents(md_text: str) -> str:
    """Extract and format table of contents from markdown text."""
    lines = md_text.split('\n')
    toc_items = []
    
    for line in lines:
        if line.startswith('#'):
            level = len(line) - len(line.lstrip('#'))
            title = line.lstrip('# ').strip()
            if title:
                indent = "  " * (level - 1)
                toc_items.append(f"{indent}- {title}")
    
    if toc_items:
        return "## 📋 Table of Contents\n\n" + "\n".join(toc_items) + "\n\n---\n\n"
    return ""

def generate_enhanced_references(topic: str) -> str:
    """Generate comprehensive references including web links, YouTube videos, and model information."""
    references_prompt = f"""
    Generate comprehensive references for the topic: {topic}
    
    Include the following categories with specific examples:
    
    1. **Academic Resources**
       - Research papers with DOI links
       - University course materials
       - Academic journals and publications
    
    2. **YouTube Educational Content**
       - Specific video recommendations with channel names
       - Educational series and playlists
       - Tutorial channels and expert content creators
    
    3. **Web Resources**
       - Official documentation and guides
       - Interactive tutorials and courses
       - Community forums and discussion platforms
    
    4. **Books and Publications**
       - Textbooks with ISBN numbers
       - E-books and online publications
       - Industry reports and white papers
    
    5. **Tools and Software**
       - Relevant software applications
       - Online tools and platforms
       - Development environments and frameworks
    
    6. **Professional Development**
       - Certification programs
       - Online courses and MOOCs
       - Professional associations and communities
    
    Format as a structured markdown list with descriptions and links where applicable.
    Focus on high-quality, authoritative sources that would be valuable for learning this topic.
    """
    
    try:
        references = call_openai(references_prompt, user_input=topic, temperature=0.3, max_output_tokens=2000)
        return references if not references.startswith("__ERROR__") else "Error generating references"
    except Exception as e:
        return f"Error generating references: {str(e)}"

def create_image_prompt_from_topic(topic: str, context: str = "") -> str:
    """Create an optimized prompt for DALL-E based on the topic and context."""
    # Clean and enhance the topic for better image generation
    clean_topic = topic.strip()
    
    # Add context if available
    if context:
        prompt = f"Educational illustration for the topic: {clean_topic}. Context: {context[:200]}. Style: clean, professional, educational diagram or illustration suitable for learning materials."
    else:
        prompt = f"Educational illustration for the topic: {clean_topic}. Style: clean, professional, educational diagram or illustration suitable for learning materials."
    
    return prompt

# -----------------------
# OpenAI wrapper (Responses API)
# -----------------------
def call_openai(instructions: str, user_input: str = "", temperature: float = 0.0, max_output_tokens: int = 1400) -> str:
    """Call Responses API. Returns text or '__ERROR__:' prefix on exception."""
    try:
        resp = client.responses.create(
            model=MODEL_NAME,
            instructions=instructions,
            input=user_input,
            temperature=temperature,
            max_output_tokens=max_output_tokens,
        )
        # Prefer high-level output_text if present
        if hasattr(resp, "output_text") and resp.output_text:
            return resp.output_text.strip()
        # Fallback: assemble from resp.output
        parts = []
        for item in getattr(resp, "output", []) or []:
            if isinstance(item, dict):
                content = item.get("content")
                if isinstance(content, list):
                    for c in content:
                        if isinstance(c, dict) and c.get("type") == "output_text":
                            parts.append(c.get("text", ""))
                elif isinstance(content, str):
                    parts.append(content)
        return "\n".join(parts).strip()
    except Exception as exc:
        return f"__ERROR__:{exc}"

# -----------------------
# Prompts
# -----------------------
OUTLINE_PROMPT = (
    "You are an expert instructor. Given a syllabus/topic, produce a concise numbered outline "
    "(4-8 top-level bullets) of what a comprehensive course/documentation should cover. Output as a plain numbered list. Topic:"
)

NOTES_PROMPT = (
    "You are an expert educator and technical writer. Produce a comprehensive, in-depth documentation-style "
    "learning module for the given Topic in Markdown. The document must include (in this order):\n\n"
    "1) **Title + Executive Summary** (2-3 paragraphs with key insights)\n"
    "2) **Prerequisites** (detailed list with explanations)\n"
    "3) **Learning Objectives** (5-10 specific, measurable goals)\n"
    "4) **Table of Contents** (auto-generated from headings)\n"
    "5) **Core Concepts** (fundamental principles with detailed explanations)\n"
    "6) **Detailed Sections** for each subtopic including:\n"
    "   - Comprehensive explanations with real-world context\n"
    "   - Step-by-step walkthroughs with examples\n"
    "   - Code snippets, diagrams descriptions, and practical applications\n"
    "   - Best practices and industry standards\n"
    "   - Performance considerations and optimization tips\n"
    "7) **Advanced Topics** (deeper dive into complex aspects)\n"
    "8) **Common Pitfalls and Misconceptions** (with explanations and how to avoid them)\n"
    "9) **Real-World Applications** (case studies and practical examples)\n"
    "10) **Study Plan** (structured learning path over 2-4 weeks)\n"
    "11) **Exercises & Projects** (15 problems with difficulty progression: Beginner, Intermediate, Advanced)\n"
    "12) **Answers Section** (detailed solutions with explanations)\n"
    "13) **Further Reading and Resources** including:\n"
    "    - Academic papers and research articles with DOI links\n"
    "    - YouTube educational channels and specific video recommendations\n"
    "    - Web resources, tutorials, and documentation\n"
    "    - Books and e-books with ISBN references\n"
    "    - Online courses and certification programs\n"
    "    - Community forums and discussion groups\n"
    "    - Tools and software related to the topic\n"
    "14) **Glossary of Key Terms** (comprehensive definitions)\n"
    "15) **Quick Reference Guide** (cheat sheet format)\n\n"
    "Output valid Markdown using headings (#, ##, ###), lists (-), code fences (```), tables, and callout boxes. "
    "Make it extremely comprehensive and detailed. Aim for 4000-6000 words total. "
    "Write in an engaging, professional tone suitable for both self-study and reference. "
    "Include practical examples, analogies, and visual descriptions for complex concepts. Topic:"
)

QUIZ_PROMPT = (
    "You are an assessment generator. Create EXACTLY 10 multiple-choice questions for the Topic in STRICT JSON format.\n"
    "Constraints:\n"
    "- Exactly 10 questions.\n"
    "- Exactly 4 options per question.\n"
    "- Include a 'difficulty' field with values 'Easy','Medium','Hard'.\n"
    "- Use this exact JSON form (answer is zero-based index):\n"
    "[\n"
    "  {\"question\":\"...\",\"options\":[\"optA\",\"optB\",\"optC\",\"optD\"], \"answer\": 0, \"difficulty\":\"Easy\"},\n"
    "  ...  (10 items total)\n"
    "]\n"
    "Split difficulties: 4 Easy, 3 Medium, 3 Hard (any order). DO NOT output any commentary outside the JSON. Topic:"
)

# -----------------------
# Robust JSON extractor + retry
# -----------------------
def extract_json(text: str):
    """Try multiple ways to parse JSON from model output."""
    if not text:
        raise ValueError("Empty text")
    if text.startswith("__ERROR__"):
        raise ValueError(text)
    # Direct
    try:
        return json.loads(text)
    except Exception:
        pass
    # Between first [ and last ]
    s = text.find("[")
    e = text.rfind("]")
    if s != -1 and e != -1 and e > s:
        chunk = text[s:e+1]
        try:
            return json.loads(chunk)
        except Exception:
            pass
    # Try JSON inside <JSON>...</JSON> tags (if model wrapped)
    m = re.search(r"<JSON>(.*)</JSON>", text, flags=re.DOTALL)
    if m:
        try:
            return json.loads(m.group(1).strip())
        except Exception:
            pass
    # Try single->double quote sanitize
    try:
        cand = text.strip().replace("\n", " ")
        cand = re.sub(r"(\w)'(\w)", r"\1’\2", cand)
        cand = cand.replace("'", '"')
        return json.loads(cand)
    except Exception as exc:
        raise ValueError("Could not parse JSON from model output: " + str(exc))

def generate_quiz_with_retries(topic: str, attempts: int = 2) -> List[Dict]:
    """Try to generate & parse quiz JSON with 1-2 attempts, using different instructions if needed."""
    # first attempt: normal strict prompt
    prompt = QUIZ_PROMPT
    for i in range(attempts):
        raw = call_openai(prompt, user_input=topic, temperature=0.0, max_output_tokens=1400)
        try:
            parsed = extract_json(raw)
            # validate structure
            if not isinstance(parsed, list) or len(parsed) != 10:
                raise ValueError("Parsed JSON not length 10")
            normalized = []
            for item in parsed:
                if not isinstance(item, dict):
                    raise ValueError("Invalid item")
                q = item.get("question")
                opts = item.get("options")
                ans = item.get("answer")
                diff = item.get("difficulty", "Medium")
                if not q or not isinstance(opts, list) or len(opts) != 4 or ans not in [0,1,2,3]:
                    raise ValueError("Invalid question format")
                normalized.append({
                    "question": q.strip(),
                    "options": [str(x).strip() for x in opts],
                    "answer": int(ans),
                    "difficulty": diff
                })
            return normalized
        except Exception as e:
            # second attempt: ask model to wrap JSON with <JSON>...</JSON> and nothing else
            if i == 0:
                prompt = "IMPORTANT: Output ONLY the JSON array. If you add comments, wrap the JSON inside <JSON> ... </JSON> tags. " + QUIZ_PROMPT
                continue
            # else fallback to raising and let caller fallback
            raise

# -----------------------
# PDF helpers (using fpdf, but transliterate unicode -> ascii with unidecode)
# -----------------------
def global_safe_multi_cell(pdf, h, txt):
    safe_txt = unidecode(str(txt))
    try:
        pdf.multi_cell(0, h, safe_txt)
    except Exception:
        # Fallback for FPDFException: Not enough horizontal space (long unbroken strings)
        chunks = [safe_txt[i:i+70] for i in range(0, len(safe_txt), 70)]
        for chunk in chunks:
            try:
                pdf.multi_cell(0, h, chunk)
            except Exception:
                pass

def markdown_to_pdf_bytes(md_text: str, title: str = "Study Notes", image_urls: List[str] = None) -> BytesIO:
    """Convert markdown to PDF with embedded images and structured formatting."""
    safe_text = unidecode(md_text)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    def safe_multi_cell(h, txt):
        global_safe_multi_cell(pdf, h, txt)

    # Title
    pdf.set_font("Arial", "B", 18)
    pdf.cell(0, 12, title, ln=True, align="C")
    pdf.ln(6)

    # Markdown Parsing (basic headers, bullets, and paragraphs)
    in_code = False
    for raw in safe_text.splitlines():
        line = raw.rstrip()
        
        # skip empty
        if line.strip() == "":
            pdf.ln(3)
            continue
            
        if line.strip().startswith("```"):
            in_code = not in_code
            if in_code:
                pdf.set_font("Courier", size=10)
            else:
                pdf.set_font("Arial", size=11)
            pdf.ln(2)
            continue

        if in_code:
            safe_multi_cell(6, line)
            continue

        # Headers
        if line.startswith("# "):
            pdf.ln(4)
            pdf.set_font("Arial", "B", 16)
            safe_multi_cell(8, line.replace("# ", "").replace("**", "").strip())
            pdf.set_font("Arial", size=11)
        elif line.startswith("## "):
            pdf.ln(3)
            pdf.set_font("Arial", "B", 14)
            safe_multi_cell(7, line.replace("## ", "").replace("**", "").strip())
            pdf.set_font("Arial", size=11)
        elif line.startswith("### "):
            pdf.ln(2)
            pdf.set_font("Arial", "B", 12)
            safe_multi_cell(6, line.replace("### ", "").replace("**", "").strip())
            pdf.set_font("Arial", size=11)
        elif line.startswith("**") and line.endswith("**"):
            pdf.set_font("Arial", "B", 11)
            safe_multi_cell(6, line[2:-2])
            pdf.set_font("Arial", size=11)
        elif re.match(r"^\s*([-*])\s+", line):
            # bullet
            pdf.set_font("Arial", size=11)
            bullet = "\u2022" + " " + re.sub(r"^\s*([-*])\s+", "", line).replace("**", "")
            # fallback to ascii bullet if unidecode messes it
            bullet_safe = bullet.encode('latin-1', 'replace').decode('latin-1')
            safe_multi_cell(6, bullet_safe)
        elif re.match(r"^\s*\d+\.\s+", line):
            pdf.set_font("Arial", size=11)
            safe_multi_cell(6, line.replace("**", "").strip())
        else:
            # text cleanup (remove inline bold formatting **)
            clean_line = line.replace("**", "")
            pdf.set_font("Arial", size=11)
            safe_multi_cell(6, clean_line)


    # Helper: download image bytes from URL or decode base64 data URI
    def _get_image_bytes(img_src):
        import base64 as b64mod
        if img_src.startswith("data:image"):
            # data:image/png;base64,XXXX
            header, encoded = img_src.split(",", 1)
            return b64mod.b64decode(encoded)
        else:
            r = requests.get(img_src, timeout=10)
            if r.status_code == 200:
                return r.content
        return None

    # Embed Images
    if image_urls:
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "Generated AI Illustrations", ln=True, align="C")
        pdf.ln(5)
        
        for i, img_url in enumerate(image_urls):
            try:
                img_bytes = _get_image_bytes(img_url)
                if img_bytes:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(img_bytes)
                        tmp_name = tmp.name
                    pdf.image(tmp_name, w=160, x=25)
                    pdf.ln(10)
                    os.unlink(tmp_name)
                else:
                    pdf.set_font("Arial", "I", 11)
                    pdf.cell(0, 6, f"[Image {i+1} failed to load]", ln=True)
            except Exception as e:
                pdf.set_font("Arial", "I", 11)
                pdf.cell(0, 6, f"[Error loading Image {i+1}: {str(e)}]", ln=True)

    out_data = pdf.output(dest="S")
    if hasattr(out_data, "encode"):
        out_data = out_data.encode("latin-1", errors="ignore")
    return BytesIO(bytes(out_data))

def markdown_to_pptx_bytes(md_text: str, title: str = "Presentation", image_urls: List[str] = None) -> BytesIO:
    """Convert markdown to a PowerPoint presentation."""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    title_shape.text = title
    subtitle.text = "Generated by Virtual AI Tutor"
    
    # Body slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    current_slide = None
    tf = None
    
    for raw in md_text.splitlines():
        line = raw.rstrip()
        if not line:
            continue
            
        clean_line = unidecode(line).replace("**", "").replace("`", "")
        
        # New slide on Header 1 or 2
        if clean_line.startswith("# ") or clean_line.startswith("## "):
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = current_slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = clean_line.replace("# ", "").replace("## ", "").strip()
            tf = body_shape.text_frame
            tf.clear() # clear default bullet
        elif current_slide and tf:
            # Bullet point or regular text
            p = tf.add_paragraph()
            if clean_line.startswith("- ") or clean_line.startswith("* "):
                p.text = clean_line[2:].strip()
                p.level = 0
            elif clean_line.startswith("### "):
                p.text = clean_line.replace("### ", "").strip()
                p.level = 0
                p.font.bold = True
            else:
                p.text = clean_line
                p.level = 1 if len(tf.paragraphs) > 1 else 0  # sub-bullet if it's not the first line
    
    # Image slides
    if image_urls:
        blank_slide_layout = prs.slide_layouts[6]
        for i, img_url in enumerate(image_urls):
            try:
                import base64 as b64mod
                if img_url.startswith("data:image"):
                    header, encoded = img_url.split(",", 1)
                    img_bytes = b64mod.b64decode(encoded)
                else:
                    resp = requests.get(img_url, timeout=10)
                    img_bytes = resp.content if resp.status_code == 200 else None
                if img_bytes:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(img_bytes)
                        tmp_name = tmp.name
                    
                    img_slide = prs.slides.add_slide(blank_slide_layout)
                    # Add simple title
                    txBox = img_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                    txBox.text_frame.text = f"Illustration {i+1}"
                    
                    # Add picture (centered roughly)
                    img_slide.shapes.add_picture(tmp_name, Inches(1.5), Inches(1.5), width=Inches(7))
                    os.unlink(tmp_name)
            except Exception:
                pass
                
    # Apply Premium Dark Theme to all slides natively
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    
    for slide in prs.slides:
        # Enforce constant dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(18, 18, 18)  # Pitch dark grey
        
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            
            tf = shape.text_frame
            # Ensure text stays within bounds
            tf.word_wrap = True 
            
            if shape == slide.shapes.title:
                # Premium Header Banner
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(99, 91, 255) # SaaS Blurple
                shape.left = 0
                shape.top = 0
                shape.width = Inches(10)
                shape.height = Inches(1.2)
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        if hasattr(run.font, "name"):
                            run.font.name = 'Arial'
                        run.font.bold = True
            else:
                # Enforce constant body positioning and typography
                shape.left = Inches(0.5)
                shape.top = Inches(1.5)
                shape.width = Inches(9)
                shape.height = Inches(5.8)
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                for paragraph in tf.paragraphs:
                    paragraph.space_after = Pt(12)
                    # paragraph.line_spacing = 1.2 # Optional
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(240, 240, 240)
                        if hasattr(run.font, "name"):
                            run.font.name = 'Inter' # Premium sans-serif
                        run.font.size = Pt(18)
    
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def quiz_pdf_bytes(quiz_list: List[Dict], title: str = "Quiz") -> BytesIO:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, title, ln=True)
    pdf.ln(6)
    pdf.set_font("Arial", size=12)
    for i, q in enumerate(quiz_list, start=1):
        global_safe_multi_cell(pdf, 8, f"Q{i}. ({q.get('difficulty','')}) {q['question']}")
        for idx, opt in enumerate(q.get("options", [])):
            global_safe_multi_cell(pdf, 8, f"   {chr(65+idx)}) {opt}")
        pdf.ln(2)
    out_data = pdf.output(dest="S")
    if hasattr(out_data, "encode"):
        out_data = out_data.encode("latin-1", errors="ignore")
    return BytesIO(bytes(out_data))

def answer_key_pdf_bytes(quiz_list: List[Dict], title: str = "Answer Key") -> BytesIO:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, title, ln=True)
    pdf.ln(6)
    pdf.set_font("Arial", size=12)
    for i, q in enumerate(quiz_list, start=1):
        ans_idx = q.get("answer")
        global_safe_multi_cell(pdf, 8, f"Q{i}. {q['question']}")
        global_safe_multi_cell(pdf, 8, f"Correct: {chr(65+ans_idx)}) {q['options'][ans_idx]}")
        pdf.ln(2)
    out_data = pdf.output(dest="S")
    if hasattr(out_data, "encode"):
        out_data = out_data.encode("latin-1", errors="ignore")
    return BytesIO(bytes(out_data))

# -----------------------
@st.cache_data(show_spinner=False)
def get_cached_pdf(md_text: str, title: str, image_urls_tuple: tuple) -> BytesIO:
    urls = list(image_urls_tuple) if image_urls_tuple else None
    return markdown_to_pdf_bytes(md_text, title, urls)

@st.cache_data(show_spinner=False)
def get_cached_pptx(md_text: str, title: str, image_urls_tuple: tuple) -> BytesIO:
    urls = list(image_urls_tuple) if image_urls_tuple else None
    return markdown_to_pptx_bytes(md_text, title, urls)

@st.cache_data(show_spinner=False)
def get_cached_quiz_pdf(quiz_json: str, title: str) -> BytesIO:
    import json
    return quiz_pdf_bytes(json.loads(quiz_json), title)

@st.cache_data(show_spinner=False)
def get_cached_ans_pdf(quiz_json: str, title: str) -> BytesIO:
    import json
    return answer_key_pdf_bytes(json.loads(quiz_json), title)

# -----------------------
# Streamlit UI (outline -> confirm -> generate)
# -----------------------

# -----------------------
# Custom CSS Theme — Premium SaaS (Vercel/Linear Style)
# -----------------------
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
@import url('https://cdn.jsdelivr.net/npm/bootstrap-icons@1.11.1/font/bootstrap-icons.css');

* { font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif; }

/* --- Base --- */
.stApp {
    background-color: #09090b;
    background-image:
        radial-gradient(ellipse 80% 50% at 50% -20%, rgba(120, 119, 198, 0.15), transparent),
        radial-gradient(ellipse 60% 40% at 80% 60%, rgba(59, 130, 246, 0.06), transparent),
        radial-gradient(ellipse 60% 40% at 20% 80%, rgba(168, 85, 247, 0.05), transparent);
    color: #e4e4e7;
}
#MainMenu, footer { visibility: hidden; }

/* Position header absolutely to remove the top gap, keeping the sidebar toggle button visible */
header[data-testid="stHeader"] {
    position: absolute !important;
    background: transparent !important;
    border-bottom: none !important;
}
header[data-testid="stHeader"] [data-testid="stHeaderActionButton"],
header[data-testid="stHeader"] [data-testid="stMainMenu"] {
    display: none !important;
}
header[data-testid="stHeader"] button {
    color: #fafafa !important;
}
.block-container {
    padding-top: 1.5rem !important;
    padding-bottom: 2rem !important;
}

/* Style Streamlit's native border container to look like our premium cards */
div[data-testid="stVerticalBlockBorderWrapper"] {
    background: rgba(24, 24, 27, 0.45) !important;
    border: 1px solid rgba(63, 63, 70, 0.4) !important;
    border-radius: 16px !important;
    padding: 1.5rem !important;
    backdrop-filter: blur(10px) !important;
    transition: border-color 0.3s ease, box-shadow 0.3s ease !important;
}
div[data-testid="stVerticalBlockBorderWrapper"]:hover {
    border-color: rgba(99, 91, 255, 0.35) !important;
    box-shadow: 0 8px 30px rgba(99, 91, 255, 0.05) !important;
}

::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: #27272a; border-radius: 10px; }
::-webkit-scrollbar-thumb:hover { background: #3f3f46; }

/* --- Sidebar --- */
section[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #0c0c0f 0%, #09090b 100%) !important;
    border-right: 1px solid rgba(63, 63, 70, 0.5) !important;
}
section[data-testid="stSidebar"] .stMarkdown h3 { font-size: 0.85rem !important; text-transform: uppercase !important; letter-spacing: 0.08em !important; color: #a1a1aa !important; }

/* --- Headings --- */
.stMarkdown h1, .stMarkdown h2, .stMarkdown h3,
.stMarkdown h4, .stMarkdown h5 { color: #fafafa !important; font-weight: 600 !important; letter-spacing: -0.025em; }
.stMarkdown { color: #a1a1aa; }

/* --- Buttons --- */
.stButton > button {
    background: linear-gradient(135deg, #ffffff 0%, #e4e4e7 100%) !important;
    color: #09090b !important;
    border: none !important;
    border-radius: 10px !important;
    font-weight: 600 !important;
    font-size: 0.875rem !important;
    padding: 0.6rem 1.25rem !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    box-shadow: 0 1px 3px rgba(0,0,0,0.3), 0 0 0 1px rgba(255,255,255,0.05) !important;
}
.stButton > button:hover {
    background: linear-gradient(135deg, #f4f4f5 0%, #d4d4d8 100%) !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 4px 12px rgba(0,0,0,0.4), 0 0 0 1px rgba(255,255,255,0.1) !important;
}
.stButton > button:active { transform: translateY(0px) !important; }
.stButton > button:disabled { background: #18181b !important; color: #52525b !important; box-shadow: none !important; }

.stDownloadButton > button {
    background: rgba(24, 24, 27, 0.8) !important;
    color: #e4e4e7 !important;
    border: 1px solid #27272a !important;
    border-radius: 10px !important;
    font-weight: 500 !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    backdrop-filter: blur(10px) !important;
}
.stDownloadButton > button:hover {
    border-color: #635BFF !important;
    background: rgba(99, 91, 255, 0.08) !important;
    box-shadow: 0 0 20px rgba(99, 91, 255, 0.15) !important;
}

/* --- Inputs --- */
.stTextInput > div > div > input {
    background: rgba(24, 24, 27, 0.6) !important;
    border: 1px solid #27272a !important;
    border-radius: 12px !important;
    color: #fafafa !important;
    padding: 0.75rem 1rem !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    font-size: 1rem !important;
    backdrop-filter: blur(10px) !important;
}
.stTextInput > div > div > input:focus {
    border-color: #635BFF !important;
    box-shadow: 0 0 0 3px rgba(99, 91, 255, 0.15), 0 0 20px rgba(99, 91, 255, 0.1) !important;
}
.stTextInput > div > div > input::placeholder { color: #52525b !important; }

/* --- Tabs --- */
div[data-testid="stTabs"] > div > div {
    background: transparent;
    border-bottom: 1px solid rgba(39, 39, 42, 0.8);
    gap: 0;
}
button[data-baseweb="tab"] {
    background: transparent !important;
    color: #71717a !important;
    font-weight: 500 !important;
    font-size: 0.875rem !important;
    border-radius: 0 !important;
    border-bottom: 2px solid transparent !important;
    padding: 0.75rem 1.25rem !important;
    transition: all 0.3s ease !important;
}
button[data-baseweb="tab"]:hover { color: #d4d4d8 !important; }
button[data-baseweb="tab"][aria-selected="true"] {
    color: #fafafa !important;
    font-weight: 600 !important;
    border-bottom: 2px solid #635BFF !important;
    box-shadow: 0 2px 8px rgba(99, 91, 255, 0.3) !important;
    background: transparent !important;
}

/* --- Progress --- */
.stProgress > div > div > div {
    background: linear-gradient(90deg, #635BFF, #a78bfa) !important;
    border-radius: 10px !important;
}

/* --- Expanders --- */
.stExpander, [data-testid="stExpander"] {
    background: rgba(24, 24, 27, 0.4) !important;
    border: 1px solid rgba(39, 39, 42, 0.6) !important;
    border-radius: 12px !important;
    backdrop-filter: blur(10px) !important;
}

/* --- Glassmorphism Card --- */
.glass-card {
    background: rgba(24, 24, 27, 0.5);
    border: 1px solid rgba(63, 63, 70, 0.4);
    border-radius: 16px;
    padding: 2rem;
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    transition: all 0.3s ease;
}
.glass-card:hover {
    border-color: rgba(99, 91, 255, 0.3);
    box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3);
}

/* --- Metric Cards --- */
.metric-card {
    background: rgba(24, 24, 27, 0.6);
    border: 1px solid rgba(63, 63, 70, 0.4);
    border-radius: 16px;
    padding: 1.75rem 1.25rem;
    text-align: center;
    backdrop-filter: blur(10px);
    transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
    position: relative;
    overflow: hidden;
}
.metric-card::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 3px;
    background: linear-gradient(90deg, #635BFF, #a78bfa, #635BFF);
    background-size: 200% 100%;
    animation: shimmer 3s ease infinite;
    opacity: 0;
    transition: opacity 0.3s ease;
}
.metric-card:hover::before { opacity: 1; }
.metric-card:hover {
    border-color: rgba(99, 91, 255, 0.3);
    transform: translateY(-4px);
    box-shadow: 0 12px 40px rgba(99, 91, 255, 0.12);
}
.metric-icon { font-size: 1.5rem; margin-bottom: 0.5rem; color: #a78bfa; }
.metric-value {
    font-size: 2.5rem;
    font-weight: 800;
    color: #fafafa;
    letter-spacing: -0.03em;
    line-height: 1;
    background: linear-gradient(135deg, #fafafa, #a1a1aa);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}
.metric-label {
    font-size: 0.7rem;
    color: #71717a;
    margin-top: 0.6rem;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.1em;
}

@keyframes shimmer {
    0% { background-position: 200% 0; }
    100% { background-position: -200% 0; }
}

/* --- Badges --- */
.badge-easy { background: rgba(5, 46, 22, 0.8); color: #34d399; border: 1px solid rgba(6, 95, 70, 0.6); padding: 3px 10px; border-radius: 999px; font-size: 0.7rem; font-weight: 600; }
.badge-medium { background: rgba(69, 26, 3, 0.8); color: #fbbf24; border: 1px solid rgba(120, 53, 15, 0.6); padding: 3px 10px; border-radius: 999px; font-size: 0.7rem; font-weight: 600; }
.badge-hard { background: rgba(69, 10, 10, 0.8); color: #f87171; border: 1px solid rgba(127, 29, 29, 0.6); padding: 3px 10px; border-radius: 999px; font-size: 0.7rem; font-weight: 600; }

/* --- Radio (Quiz) --- */
.stRadio > div { gap: 0.4rem !important; }
.stRadio > div > label {
    background: rgba(24, 24, 27, 0.5) !important;
    border: 1px solid rgba(39, 39, 42, 0.6) !important;
    border-radius: 10px !important;
    padding: 0.6rem 1rem !important;
    transition: all 0.3s ease !important;
    cursor: pointer !important;
    backdrop-filter: blur(5px) !important;
}
.stRadio > div > label:hover {
    border-color: rgba(99, 91, 255, 0.4) !important;
    background: rgba(99, 91, 255, 0.05) !important;
}

/* --- Code --- */
.stMarkdown code { background: rgba(24, 24, 27, 0.8) !important; color: #e4e4e7 !important; border: 1px solid #27272a !important; border-radius: 6px !important; padding: 0.15em 0.4em !important; font-size: 0.85em !important; }
.stMarkdown pre { background: rgba(9, 9, 11, 0.9) !important; border: 1px solid #27272a !important; border-radius: 12px !important; }

/* --- Slider --- */
.stSlider > div > div > div > div { background: #635BFF !important; }

/* --- Dividers --- */
hr { border-color: rgba(39, 39, 42, 0.5) !important; }

/* --- Alerts --- */
div[data-testid="stAlert"] { border-radius: 12px !important; border: 1px solid rgba(39, 39, 42, 0.6) !important; background: rgba(24, 24, 27, 0.6) !important; color: #e4e4e7 !important; backdrop-filter: blur(10px) !important; }

/* --- Footer --- */
.custom-footer {
    text-align: center;
    padding: 2.5rem 0 1.5rem;
    color: #52525b;
    font-size: 0.8rem;
    border-top: 1px solid rgba(39, 39, 42, 0.5);
    margin-top: 4rem;
    letter-spacing: 0.02em;
}
.custom-footer a { color: #a78bfa; text-decoration: none; transition: color 0.2s ease; }
.custom-footer a:hover { color: #635BFF; }

/* --- Hero --- */
.sf-header {
    display: flex;
    align-items: center;
    justify-content: center;
    flex-direction: column;
    gap: 16px;
    margin: 0.5rem 0 1rem 0;
    text-align: center;
    position: relative;
}
.sf-header-badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    font-size: 0.75rem;
    font-weight: 600;
    color: #a78bfa;
    background: rgba(99, 91, 255, 0.1);
    border: 1px solid rgba(99, 91, 255, 0.2);
    padding: 6px 16px;
    border-radius: 999px;
    margin-bottom: 0.5rem;
    animation: fadeInDown 0.6s ease;
}
.sf-header-title {
    font-size: 3.5rem;
    font-weight: 900;
    letter-spacing: -0.05em;
    margin: 0;
    line-height: 1.05;
    background: linear-gradient(135deg, #fafafa 0%, #a1a1aa 50%, #fafafa 100%);
    background-size: 200% 200%;
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    animation: gradientShift 6s ease infinite;
}
.sf-desc {
    font-size: 1.1rem;
    color: #71717a;
    line-height: 1.6;
    max-width: 550px;
    margin: 0 auto;
    text-align: center;
    font-weight: 400;
}

@keyframes fadeInDown {
    from { opacity: 0; transform: translateY(-10px); }
    to { opacity: 1; transform: translateY(0); }
}
@keyframes gradientShift {
    0%, 100% { background-position: 0% 50%; }
    50% { background-position: 100% 50%; }
}

/* --- Section Card --- */
.section-card {
    background: rgba(24, 24, 27, 0.4);
    border: 1px solid rgba(63, 63, 70, 0.3);
    border-radius: 16px;
    padding: 1.75rem;
    margin: 1rem 0;
    backdrop-filter: blur(10px);
}
.section-card-header {
    font-size: 0.7rem;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    color: #635BFF;
    margin-bottom: 0.75rem;
}

/* --- Download Card --- */
.dl-card {
    background: rgba(24, 24, 27, 0.4);
    border: 1px solid rgba(63, 63, 70, 0.3);
    border-radius: 14px;
    padding: 1.25rem;
    text-align: center;
    transition: all 0.3s ease;
}
.dl-card:hover {
    border-color: rgba(99, 91, 255, 0.3);
    background: rgba(99, 91, 255, 0.04);
}
.dl-icon { font-size: 1.8rem; margin-bottom: 0.4rem; }
.dl-label { font-size: 0.75rem; color: #71717a; font-weight: 500; margin-top: 0.3rem; }

/* --- Outline Card --- */
.outline-card {
    background: rgba(24, 24, 27, 0.4);
    border: 1px solid rgba(63, 63, 70, 0.3);
    border-radius: 16px;
    padding: 1.5rem;
    backdrop-filter: blur(10px);
}
.outline-card ol, .outline-card ul {
    color: #d4d4d8;
    line-height: 1.8;
}

/* --- Score Card --- */
.score-card {
    background: rgba(24, 24, 27, 0.5);
    border: 1px solid rgba(63, 63, 70, 0.3);
    border-radius: 20px;
    padding: 2rem;
    text-align: center;
    margin: 1rem 0;
    backdrop-filter: blur(10px);
    position: relative;
    overflow: hidden;
}
.score-card::before {
    content: '';
    position: absolute;
    top: -50%; left: -50%;
    width: 200%; height: 200%;
    background: radial-gradient(circle, rgba(99, 91, 255, 0.05) 0%, transparent 70%);
    animation: pulse 4s ease infinite;
}
@keyframes pulse {
    0%, 100% { transform: scale(1); opacity: 0.5; }
    50% { transform: scale(1.05); opacity: 1; }
}

/* --- Feature Pills --- */
.feature-pill {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    background: rgba(24, 24, 27, 0.6);
    border: 1px solid rgba(63, 63, 70, 0.4);
    border-radius: 999px;
    padding: 8px 16px;
    font-size: 0.8rem;
    color: #a1a1aa;
    margin: 4px;
    transition: all 0.2s ease;
}
.feature-pill:hover {
    border-color: rgba(99, 91, 255, 0.4);
    color: #d4d4d8;
}
</style>
""", unsafe_allow_html=True)

# --- Header ---
st.markdown("""
<div class="sf-header">
    <div class="sf-header-badge"><i class="bi bi-stars"></i> Virtual AI Tutor 1.0</div>
    <h1 class="sf-header-title">Master any topic,<br>instantly.</h1>
    <p class="sf-desc">Generate comprehensive study notes, interactive quizzes, and AI-powered illustrations from a single prompt.</p>
    <div style="text-align:center; margin-top:1rem;">
        <span class="feature-pill"><i class="bi bi-journal-text"></i> Study Notes</span>
        <span class="feature-pill"><i class="bi bi-image"></i> AI Images</span>
        <span class="feature-pill"><i class="bi bi-patch-question"></i> Quizzes</span>
        <span class="feature-pill"><i class="bi bi-journal-bookmark"></i> References</span>
        <span class="feature-pill"><i class="bi bi-file-earmark-arrow-down"></i> PDF/PPTX Export</span>
    </div>
</div>
<div style="height: 1rem;"></div>
""", unsafe_allow_html=True)



with st.sidebar:
    st.markdown("### Virtual AI Tutor")
    st.caption("Intelligent Learning Assistant")
    st.markdown("---")
    
    st.markdown('### <i class="bi bi-clock-history"></i> History', unsafe_allow_html=True)
    history_items = load_history()
    if history_items:
        with st.expander("Previously Generated Topics"):
            for idx, item in enumerate(history_items):
                label = f"{item.get('topic', 'Unknown')} ({item.get('timestamp', '')[:10]})"
                if st.button(label, key=f"hist_{idx}"):
                    st.session_state['topic'] = item.get('topic')
                    st.session_state['outline'] = item.get('outline')
                    st.session_state['notes_md'] = item.get('notes_md')
                    st.session_state['generated_images'] = item.get('generated_images', [])
                    
                    loaded_quiz = item.get('quiz', [])
                    st.session_state['quiz'] = loaded_quiz
                    st.session_state['answers'] = [None] * len(loaded_quiz)
                    st.session_state['show_key'] = False
                    st.session_state['quiz_submitted'] = False
                    
                    st.session_state['enhanced_references'] = item.get('enhanced_references', '')
                    st.rerun()
    else:
        st.caption("No history yet.")
    st.markdown("---")
    st.markdown("**Getting started**")
    st.markdown("""
1. Enter a topic below  
2. Review the generated outline  
3. Confirm to generate all materials  
4. Download as PDF
""")
    st.markdown(f"""
<span style="display:inline-block; background:rgba(99,91,255,0.1); border:1px solid rgba(99,91,255,0.2); border-radius:4px; padding:2px 8px; font-size:0.75rem; font-weight:600; color:#635BFF;">
Model: {MODEL_NAME}
</span>
""", unsafe_allow_html=True)
    
    st.markdown("---")
    st.markdown("**Settings**")
    num_images = st.slider("Number of images", min_value=1, max_value=5, value=3)
    st.session_state['num_images'] = num_images
    
    st.markdown("---")
    st.markdown("**Quick Tools**")
    
    with st.expander("Image Generation"):
        image_topic = st.text_input("Topic", placeholder="e.g., Machine Learning", key="quick_image")
        if st.button("Generate", key="btn_image"):
            if image_topic.strip():
                with st.spinner("Generating..."):
                    image_prompt = create_image_prompt_from_topic(image_topic)
                    image_url = generate_image_with_dalle(image_prompt)
                    if not image_url.startswith("Error"):
                        st.image(image_url, caption=f"Generated for: {image_topic}", width="stretch")
                    else:
                        st.error(f"Failed: {image_url}")
            else:
                st.warning("Enter a topic first.")
    
    with st.expander("References"):
        ref_topic = st.text_input("Topic", placeholder="e.g., Machine Learning", key="quick_refs")
        if st.button("Generate", key="btn_refs"):
            if ref_topic.strip():
                with st.spinner("Generating..."):
                    refs_result = generate_enhanced_references(ref_topic)
                    if not refs_result.startswith("Error"):
                        st.markdown(refs_result)
                    else:
                        st.error(f"Failed: {refs_result}")
            else:
                st.warning("Enter a topic first.")



with st.container(border=True):
    st.markdown('<div class="section-card-header"><i class="bi bi-pencil-square"></i> Enter Your Topic</div>', unsafe_allow_html=True)
    topic = st.text_input("Topic", value=st.session_state.get('topic', ''), placeholder="e.g., Machine Learning, Data Structures, Quantum Physics...", label_visibility="collapsed")

col_btn1, col_btn2 = st.columns([1.5, 4.5])
expand_clicked = False
with col_btn1:
    expand_clicked = st.button("Expand Outline", type="primary", use_container_width=True)

if expand_clicked:
    if not topic.strip():
        st.warning("Please enter a topic first.")
    else:
        with st.spinner("Generating outline..."):
            out_raw = call_openai(OUTLINE_PROMPT, user_input=topic, temperature=0.0, max_output_tokens=300)
        if out_raw.startswith("__ERROR__"):
            st.error("Error generating outline: " + out_raw)
        else:
            st.session_state['outline'] = out_raw
            st.session_state['outline_success'] = True

if st.session_state.get('outline_success') and 'outline' in st.session_state and 'notes_md' not in st.session_state:
    st.success("Outline ready! Review and confirm to generate content.")

# show outline and confirm
if 'outline' in st.session_state:
    st.markdown('')
    col_outline, col_action = st.columns([2, 1])
    with col_outline:
        with st.container(border=True):
            st.markdown('<div class="section-card-header"><i class="bi bi-list-task"></i> Generated Outline</div>', unsafe_allow_html=True)
            st.markdown(st.session_state['outline'])
    with col_action:
        with st.container(border=True):
            st.markdown('<div class="section-card-header" style="text-align:center;"><i class="bi bi-rocket-takeoff"></i> Ready to Generate</div>', unsafe_allow_html=True)
            st.markdown('<p style="text-align:center; color:#71717a; font-size:0.875rem; margin-bottom:1rem;">This will generate notes, images, quizzes, and references.</p>', unsafe_allow_html=True)
            if st.button("Confirm & Generate", type="primary", use_container_width=True):
                # Progress tracking
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # Generate notes
                status_text.text("Generating notes...")
                progress_bar.progress(20)
                notes_md = call_openai(NOTES_PROMPT, user_input=topic, temperature=0.0, max_output_tokens=6000)
                if notes_md.startswith("__ERROR__"):
                    st.error("Error generating notes: " + notes_md)
                else:
                    st.session_state['notes_md'] = notes_md
                    progress_bar.progress(40)
                    
                # Generate multiple images with DALL-E
                status_text.text("Generating images...")
                progress_bar.progress(60)
                num_images = st.session_state.get('num_images', 3)
                image_urls = generate_multiple_images_with_dalle(topic, num_images, notes_md[:500])
                if image_urls:
                    st.session_state['generated_images'] = image_urls
                    st.session_state['image_prompts'] = [create_image_prompt_from_topic(f"Aspect {i+1} of {topic}", notes_md[:200]) for i in range(len(image_urls))]
                    progress_bar.progress(80)
                else:
                    st.warning("Failed to generate images.")
                
                # Generate enhanced references
                status_text.text("Generating references...")
                enhanced_refs = generate_enhanced_references(topic)
                if not enhanced_refs.startswith("Error"):
                    st.session_state['enhanced_references'] = enhanced_refs
                    progress_bar.progress(90)
                else:
                    st.warning("Failed to generate references.")

                # Generate quiz
                status_text.text("Generating quiz...")
                progress_bar.progress(95)
                try:
                    quiz_list = generate_quiz_with_retries(topic, attempts=2)
                    st.session_state['quiz'] = quiz_list
                    st.session_state['answers'] = [None] * len(quiz_list)
                    st.session_state['show_key'] = False
                    progress_bar.progress(100)
                    status_text.text("Done.")
                    st.success("All content generated successfully.")
                except Exception as e:
                    st.error("Quiz generation failed: " + str(e))
                    st.info("Using fallback quiz.")
                    # deterministic fallback
                    base = st.session_state.get('notes_md', topic)
                    lines = [ln.strip() for ln in base.splitlines() if ln.strip()]
                    fallback = []
                    for i in range(10):
                      snippet = lines[i % len(lines)] if lines else f"Fact about {topic}"
                      qtext = f"Which statement about the topic is true? ({snippet[:80]})"
                      opts = [f"{snippet} (true)", "Incorrect option A", "Incorrect option B", "Incorrect option C"]
                      diff = "Easy" if i < 4 else ("Medium" if i < 7 else "Hard")
                      fallback.append({"question": qtext, "options": opts, "answer": 0, "difficulty": diff})
                    st.session_state['quiz'] = fallback
                    st.session_state['answers'] = [None]*len(fallback)
                    st.session_state['show_key'] = False
                    progress_bar.progress(100)
                    status_text.text("Done.")
                    st.success("Content generated with fallback quiz.")
                
                # Save strictly if notes generated
                if 'notes_md' in st.session_state:
                    import datetime
                    record = {
                        "timestamp": datetime.datetime.now().isoformat(),
                        "topic": topic,
                        "outline": st.session_state.get('outline', ''),
                        "notes_md": st.session_state.get('notes_md', ''),
                        "generated_images": st.session_state.get('generated_images', []),
                        "quiz": st.session_state.get('quiz', []),
                        "enhanced_references": st.session_state.get('enhanced_references', '')
                    }
                    save_history(record)

# ===========================
# TABBED CONTENT DISPLAY
# ===========================
has_content = 'notes_md' in st.session_state or 'quiz' in st.session_state

if has_content:
    st.markdown('')

    # --- Metrics Dashboard ---
    notes_text = st.session_state.get('notes_md', '')
    quiz_data = st.session_state.get('quiz', [])
    images_list = st.session_state.get('generated_images', [])
    word_count = len(notes_text.split()) if notes_text else 0

    m1, m2, m3, m4 = st.columns(4)
    with m1:
        st.markdown(f'<div class="metric-card"><div class="metric-icon"><i class="bi bi-file-earmark-richtext"></i></div><div class="metric-value">{word_count:,}</div><div class="metric-label">Words Generated</div></div>', unsafe_allow_html=True)
    with m2:
        st.markdown(f'<div class="metric-card"><div class="metric-icon"><i class="bi bi-question-square"></i></div><div class="metric-value">{len(quiz_data)}</div><div class="metric-label">Quiz Questions</div></div>', unsafe_allow_html=True)
    with m3:
        st.markdown(f'<div class="metric-card"><div class="metric-icon"><i class="bi bi-images"></i></div><div class="metric-value">{len(images_list)}</div><div class="metric-label">Images Created</div></div>', unsafe_allow_html=True)
    with m4:
        refs_available = "✓" if st.session_state.get('enhanced_references') else "—"
        st.markdown(f'<div class="metric-card"><div class="metric-icon"><i class="bi bi-journal-bookmark"></i></div><div class="metric-value">{refs_available}</div><div class="metric-label">References</div></div>', unsafe_allow_html=True)

    st.markdown('')

    # --- Tabs ---
    tab_notes, tab_images, tab_quiz, tab_refs, tab_downloads = st.tabs([
        "Study Notes", "Images", "Quiz", "References", "Downloads"
    ])

    # ===== TAB: NOTES =====
    with tab_notes:
        if 'notes_md' in st.session_state:
            toc = extract_table_of_contents(st.session_state['notes_md'])
            if toc:
                with st.expander("Table of Contents", expanded=False):
                    st.markdown(toc)
            st.markdown(st.session_state['notes_md'])
        else:
            st.info("Notes will appear here after generation.")

    # ===== TAB: IMAGES =====
    with tab_images:
        if 'generated_images' in st.session_state and st.session_state['generated_images']:
            images = st.session_state['generated_images']
            cols = st.columns(min(len(images), 3))
            for i, image_url in enumerate(images):
                with cols[i % 3]:
                    st.image(image_url, caption=f"Illustration {i+1}", width="stretch")

            with st.expander("View image generation prompts"):
                prompts = st.session_state.get('image_prompts', [])
                for i, prompt in enumerate(prompts):
                    st.caption(f"**Image {i+1}:** {prompt}")

            if st.button("Regenerate All Images", key="regen_images"):
                with st.spinner("Generating new images..."):
                    num_imgs = st.session_state.get('num_images', 3)
                    new_urls = generate_multiple_images_with_dalle(topic, num_imgs, st.session_state.get('notes_md', '')[:500])
                    if new_urls:
                        st.session_state['generated_images'] = new_urls
                        st.success(f"Generated {len(new_urls)} new images.")
                        st.rerun()
                    else:
                        st.error("Failed to regenerate images.")
        else:
            st.info("Images will appear here after generation.")

    # ===== TAB: QUIZ =====
    with tab_quiz:
        if 'quiz' in st.session_state:
            quiz = st.session_state['quiz']

            # Quiz progress
            answers = st.session_state.get('answers', [None]*len(quiz))
            answered = sum(1 for a in answers if a is not None)
            quiz_progress = int((answered / len(quiz)) * 100) if len(quiz) else 0

            st.markdown(f"""
            <div style="background:#0e0e16; border:1px solid #1e1e2e; border-radius:8px; padding:0.85rem 1rem; margin-bottom:0.75rem;">
                <div style="display:flex; justify-content:space-between; align-items:center;">
                    <span style="font-weight:500; color:#eeeef2; font-size:0.9rem;">Progress: {answered}/{len(quiz)} answered</span>
                    <span style="color:#666680; font-size:0.85rem;">{quiz_progress}%</span>
                </div>
            </div>
            """, unsafe_allow_html=True)
            st.progress(quiz_progress)

            order = ["Easy", "Medium", "Hard"]
            badge_map = {"Easy": "badge-easy", "Medium": "badge-medium", "Hard": "badge-hard"}
            grouped = {k: [] for k in order}
            for idx, q in enumerate(quiz):
                diff = q.get('difficulty', 'Medium')
                grouped.setdefault(diff, []).append(idx)

            for diff in order:
                idxs = grouped.get(diff, [])
                if not idxs:
                    continue
                badge_class = badge_map.get(diff, "badge-medium")
                with st.expander(f"{diff} — {len(idxs)} Questions", expanded=(diff == "Easy")):
                    for idx in idxs:
                        q = quiz[idx]
                        st.markdown(f'**Q{idx+1}.** <span class="{badge_class}">{diff}</span> {q["question"]}', unsafe_allow_html=True)
                        opts_with_placeholder = ["— Select an option —"] + q['options']
                        key = f"q_{idx}"
                        prev = st.session_state['answers'][idx]
                        default_index = (prev + 1) if prev is not None else 0
                        if default_index < 0 or default_index >= len(opts_with_placeholder):
                            default_index = 0
                        selected = st.radio("Select:", opts_with_placeholder, key=key, index=default_index, label_visibility="collapsed")
                        if selected == "— Select an option —":
                            st.session_state['answers'][idx] = None
                        else:
                            try:
                                sel_idx = q['options'].index(selected)
                            except ValueError:
                                sel_idx = opts_with_placeholder.index(selected) - 1
                            st.session_state['answers'][idx] = int(sel_idx)
                        st.markdown("---")

            # Submit / Answer Key buttons
            col1, col2, col3 = st.columns([1, 1, 1])
            with col1:
                all_ans = all(a is not None for a in st.session_state.get('answers', []))
                if st.button("Submit Quiz", disabled=not all_ans, type="primary", use_container_width=True):
                    st.session_state['quiz_submitted'] = True
                    st.session_state['show_key'] = True
                    st.rerun()
            with col2:
                if st.button("View Answer Key", use_container_width=True):
                    st.session_state['show_key'] = True
            with col3:
                if st.session_state.get('quiz_submitted') or st.session_state.get('show_key'):
                    if st.button("Reset Quiz", use_container_width=True):
                        st.session_state['answers'] = [None] * len(quiz)
                        st.session_state['show_key'] = False
                        st.session_state['quiz_submitted'] = False
                        st.rerun()

            if not all_ans and not st.session_state.get('show_key'):
                st.info("Answer all questions to submit the quiz.")

            # Results display
            if st.session_state.get('show_key'):
                answers = st.session_state['answers']
                score = sum(1 for i, q in enumerate(quiz) if answers[i] == q['answer'])
                percentage = score / len(quiz) * 100

                # Score card
                score_color = "#34d399" if percentage >= 70 else ("#fbbf24" if percentage >= 50 else "#f87171")
                st.markdown(f"""
                <div style="background:#0e0e16; border:1px solid #1e1e2e; border-radius:10px; padding:1.5rem; text-align:center; margin:1rem 0;">
                    <div style="font-size:2.5rem; font-weight:700; color:{score_color};">{score}/{len(quiz)}</div>
                    <div style="font-size:1.2rem; color:{score_color}; font-weight:600; margin-top:0.25rem;">{percentage:.0f}%</div>
                    <div style="color:#666680; margin-top:0.5rem; font-size:0.85rem;">
                        {"Quiz Submitted" if st.session_state.get('quiz_submitted') else "Answer Key Preview"}
                    </div>
                </div>
                """, unsafe_allow_html=True)

                if percentage >= 90:
                    st.success("Excellent — you've mastered this topic.")
                elif percentage >= 80:
                    st.success("Great work. Solid understanding.")
                elif percentage >= 70:
                    st.warning("Good effort. Review incorrect answers.")
                elif percentage >= 60:
                    st.warning("Keep studying. Focus on missed areas.")
                else:
                    st.error("More study needed. Review and try again.")

                st.markdown("**Answer Key**")
                for i, q in enumerate(quiz):
                    user_idx = answers[i]
                    correct_idx = q['answer']
                    is_correct = user_idx == correct_idx
                    icon = "✅" if is_correct else "❌"
                    user_label = chr(65 + user_idx) if user_idx is not None else "N/A"
                    correct_label = chr(65 + correct_idx)
                    st.code(f"{icon} Q{i+1}. ({q.get('difficulty','')}) {q['question']}\nCorrect: {correct_label}) {q['options'][correct_idx]}\nYour answer: {user_label}) {q['options'][user_idx] if user_idx is not None else ''}")
        else:
            st.info("Quiz will appear here after generation.")

    # ===== TAB: REFERENCES =====
    with tab_refs:
        if 'enhanced_references' in st.session_state and st.session_state['enhanced_references']:
            st.markdown(st.session_state['enhanced_references'])
            if st.button("Regenerate References", key="regen_refs"):
                with st.spinner("Generating new references..."):
                    new_refs = generate_enhanced_references(topic)
                    if not new_refs.startswith("Error"):
                        st.session_state['enhanced_references'] = new_refs
                        st.success("References regenerated.")
                        st.rerun()
                    else:
                        st.error("Failed to regenerate references.")
        else:
            st.info("References will appear here after generation.")

    # ===== TAB: DOWNLOADS =====
    with tab_downloads:
        with st.container(border=True):
            st.markdown('<div class="section-card-header"><i class="bi bi-download"></i> Export Materials</div>', unsafe_allow_html=True)
            st.markdown('<p style="color:#71717a; font-size:0.875rem; margin-bottom:1rem;">Download your generated content as PDF or PPTX files.</p>', unsafe_allow_html=True)

            dl1, dl2, dl3, dl4 = st.columns(4)
            with dl1:
                if 'notes_md' in st.session_state:
                    urls_tuple = tuple(st.session_state.get('generated_images', []))
                    notes_pdf = get_cached_pdf(st.session_state['notes_md'], f"Notes: {topic}", urls_tuple)
                    st.download_button("Study Notes (PDF)", notes_pdf, file_name="notes.pdf", mime="application/pdf", use_container_width=True)
                else:
                    st.button("Study Notes (PDF)", disabled=True, use_container_width=True)
            with dl2:
                if 'notes_md' in st.session_state:
                    urls_tuple = tuple(st.session_state.get('generated_images', []))
                    notes_pptx = get_cached_pptx(st.session_state['notes_md'], f"Presentation: {topic}", urls_tuple)
                    st.download_button("Presentation (PPTX)", notes_pptx, file_name="presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                else:
                    st.button("Presentation (PPTX)", disabled=True, use_container_width=True)
            with dl3:
                if 'quiz' in st.session_state:
                    quiz_json = json.dumps(st.session_state['quiz'])
                    quiz_pdf = get_cached_quiz_pdf(quiz_json, f"Quiz: {topic}")
                    st.download_button("Quiz (PDF)", quiz_pdf, file_name="quiz.pdf", mime="application/pdf", use_container_width=True)
                else:
                    st.button("Quiz (PDF)", disabled=True, use_container_width=True)
            with dl4:
                if 'quiz' in st.session_state:
                    quiz_json = json.dumps(st.session_state['quiz'])
                    ans_pdf = get_cached_ans_pdf(quiz_json, f"Answer Key: {topic}")
                    st.download_button("Answer Key (PDF)", ans_pdf, file_name="answer_key.pdf", mime="application/pdf", use_container_width=True)
                else:
                    st.button("Answer Key (PDF)", disabled=True, use_container_width=True)

# --- Footer ---
st.markdown("""
<div class="custom-footer">
    <div style="margin-bottom: 0.5rem;"><i class="bi bi-stars"></i> Virtual AI Tutor</div>
    Built with Streamlit & OpenAI &mdash; Powered by GPT-4o & GPT-image-1
</div>
""", unsafe_allow_html=True)
